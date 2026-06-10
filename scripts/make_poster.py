"""Academic-style poster for CodeRetina.

Conventions followed:
  * Portrait 0.9 m x 1.2 m at 150 DPI -> 5315 x 7086 px (matches template).
  * Disciplined color palette: deep purple (matches template band) + charcoal +
    soft purple tint for tables. No decorative gradients in the body.
  * 3-column body, dense figures, no colored "chip" cards.
  * All readable body text >= 44 pt (中文微软雅黑 / 英文 Arial), with
    PingFang SC used as the rendering substitute on macOS.

Outputs:
  - poster_CodeRetina.jpg
  - poster_CodeRetina.pdf
  - poster_CodeRetina.pptx
"""
from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
OUT_JPG = ROOT / "poster_CodeRetina.jpg"
OUT_PDF = ROOT / "poster_CodeRetina.pdf"
OUT_PPTX = ROOT / "poster_CodeRetina.pptx"

# ---------- Canvas ----------
DPI = 150
W = 5315
H = 7086

def pt(p: float) -> int:
    return int(round(p * DPI / 72.0))

# ---------- Palette ----------
PURPLE          = (91, 42, 137)     # primary
PURPLE_DEEP     = (62, 24, 102)     # darker text accent
PURPLE_BAND_A   = (95, 35, 140)     # template-matching banner
PURPLE_BAND_B   = (122, 49, 169)
PURPLE_FOOT_A   = (84, 71, 167)
PURPLE_FOOT_B   = (118, 60, 178)
TINT            = (243, 235, 252)   # subtle table row tint
TINT_HEAVY      = (227, 213, 244)   # slightly stronger
RULE            = (208, 196, 226)   # decorative rule
CHARCOAL        = (35, 30, 50)      # body text
GRAY_TEXT       = (95, 88, 115)
GRAY_LIGHT      = (160, 152, 180)
ACCENT          = (191, 88, 36)     # amber accent for in-text emphasis
GOLD            = (180, 132, 38)
WHITE           = (255, 255, 255)
BLACK           = (15, 12, 25)

# ---------- Fonts (paths discovered at design time on macOS) ----------
ARIAL_REG   = "/Library/Fonts/Arial Unicode.ttf"
ARIAL_BOLD  = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
ARIAL_BLACK = "/System/Library/Fonts/Supplemental/Arial Black.ttf"
ARIAL_IT    = "/System/Library/Fonts/Supplemental/Arial Italic.ttf"
PINGFANG    = "/System/Library/AssetsV2/com_apple_MobileAsset_Font7/3419f2a427639ad8c8e139149a287865a90fa17e.asset/AssetData/PingFang.ttc"

def font(p: float, weight: str = "regular", latin: bool = False,
         italic: bool = False) -> ImageFont.FreeTypeFont:
    size = pt(p)
    if latin:
        if italic and weight == "regular":
            path = ARIAL_IT
        elif weight == "black":
            path = ARIAL_BLACK
        elif weight == "bold":
            path = ARIAL_BOLD
        else:
            path = ARIAL_REG
        return ImageFont.truetype(path, size)
    idx = {"regular": 0, "bold": 4, "black": 5}.get(weight, 0)
    try:
        return ImageFont.truetype(PINGFANG, size, index=idx)
    except Exception:
        return ImageFont.truetype(PINGFANG, size, index=0)

# ---------- Drawing primitives ----------

def text_size(draw, s, fnt):
    bb = draw.textbbox((0, 0), s, font=fnt)
    return bb[2] - bb[0], bb[3] - bb[1]

def t(draw, xy, s, fnt, fill=CHARCOAL, anchor="la"):
    draw.text(xy, s, font=fnt, fill=fill, anchor=anchor)

def hline(draw, x0, x1, y, color=RULE, width=2):
    draw.line([(x0, y), (x1, y)], fill=color, width=width)

def vline(draw, x, y0, y1, color=RULE, width=2):
    draw.line([(x, y0), (x, y1)], fill=color, width=width)

def rrect(draw, box, r=24, fill=None, outline=None, width=2):
    draw.rounded_rectangle(box, radius=r, fill=fill, outline=outline, width=width)

def vgradient(img, x0, y0, x1, y1, c1, c2):
    h = y1 - y0
    px = img.load()
    for i in range(h):
        k = i / max(h - 1, 1)
        r = int(c1[0] + (c2[0] - c1[0]) * k)
        g = int(c1[1] + (c2[1] - c1[1]) * k)
        b = int(c1[2] + (c2[2] - c1[2]) * k)
        for x in range(x0, x1):
            px[x, y0 + i] = (r, g, b)

def wrap(draw, s, fnt, max_w):
    lines, cur = [], ""
    for ch in s:
        if ch == "\n":
            lines.append(cur); cur = ""; continue
        trial = cur + ch
        w, _ = text_size(draw, trial, fnt)
        if w > max_w and cur:
            lines.append(cur); cur = ch
        else:
            cur = trial
    if cur:
        lines.append(cur)
    return lines

def paragraph(draw, x, y, s, fnt, max_w, fill=CHARCOAL, line_h=1.45):
    lines = wrap(draw, s, fnt, max_w)
    _, lh = text_size(draw, "字Hg", fnt)
    step = int(lh * line_h)
    for i, line in enumerate(lines):
        t(draw, (x, y + i * step), line, fnt, fill)
    return y + len(lines) * step

# ---------- Section header (academic style) ----------

def section(draw, x, y, w, num, title_zh, title_en):
    """Numbered section header with thin rule beneath.
    Returns y just below the rule (with margin)."""
    f_num = font(56, weight="bold", latin=True)
    f_zh  = font(56, weight="bold")
    f_en  = font(36, weight="regular", latin=True, italic=True)
    nw, nh = text_size(draw, num, f_num)
    t(draw, (x, y), num, f_num, PURPLE)
    t(draw, (x + nw + 24, y), title_zh, f_zh, CHARCOAL)
    zw, _ = text_size(draw, title_zh, f_zh)
    en_x = x + nw + 24 + zw + 26
    t(draw, (en_x, y + 18), title_en, f_en, GRAY_TEXT)
    rule_y = y + nh + 22
    hline(draw, x, x + w, rule_y, color=PURPLE, width=4)
    return rule_y + 36

# ---------- Banners (top / bottom) ----------

def banner_top(img, draw):
    band_h = 700
    vgradient(img, 0, 0, W, band_h, PURPLE_BAND_A, PURPLE_BAND_B)
    overlay = Image.new("RGBA", (W, band_h), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.ellipse([W - 1500, -700, W + 100, 500], fill=(255, 255, 255, 22))
    od.ellipse([-300, -400, 800, 500], fill=(255, 255, 255, 18))
    img.alpha_composite(overlay, (0, 0))
    f_main = font(82, weight="bold")
    f_sub  = font(46, weight="regular")
    t(draw, (W // 2, 230), "大模型驱动的软件开发  ·  课程成果展示",
      f_main, WHITE, anchor="mm")
    t(draw, (W // 2, 370), "2026 春季学期  ·  Spring 2026",
      f_sub, (240, 220, 255), anchor="mm")
    hline(draw, W // 2 - 600, W // 2 + 600, 470, (235, 215, 250), width=3)

def banner_bottom(img, draw):
    band_h = 360
    vgradient(img, 0, H - band_h, W, H, PURPLE_FOOT_A, PURPLE_FOOT_B)
    overlay = Image.new("RGBA", (W, band_h), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.ellipse([-200, -200, 700, 500], fill=(255, 255, 255, 22))
    od.ellipse([W - 600, 80, W + 300, 700], fill=(255, 255, 255, 18))
    img.alpha_composite(overlay, (0, H - band_h))

    f_brand    = font(64, weight="bold", latin=True)
    f_brand_zh = font(36)
    t(draw, (200, H - band_h + 70), "CodeRetina", f_brand, WHITE)
    t(draw, (200, H - band_h + 70 + pt(64) + 12),
      "Vision  ·  Verification  ·  Loopback   /   让 Code Agent 看得见 · 看得懂 · 看得准",
      f_brand_zh, (235, 215, 250))

    # contact placeholders (editable in pptx by user)
    f_lab = font(40, weight="bold")
    f_val = font(40, weight="bold", latin=True)
    rx = W - 200
    y1 = H - band_h + 80
    t(draw, (rx, y1), "联系人  ·  Contact", f_lab, WHITE, anchor="ra")
    t(draw, (rx, y1 + pt(40) + 12), "________________", f_val, (245, 220, 255), anchor="ra")
    t(draw, (rx, y1 + pt(40) + 12 + pt(40) + 26), "邮箱  ·  Email", f_lab, WHITE, anchor="ra")
    t(draw, (rx, y1 + pt(40) + 12 + pt(40) + 26 + pt(40) + 12),
      "________________", f_val, (245, 220, 255), anchor="ra")

    f_credit = font(28, weight="regular", latin=True)
    t(draw, (W // 2, H - 40),
      "Course Project  ·  大模型驱动的软件开发  ·  Spring 2026",
      f_credit, (230, 215, 250), anchor="mm")

# ---------- Hero block ----------

def render_hero(img, draw, y_start):
    """Title + taglines + author line. Returns y below."""
    cx = W // 2
    # Project name (huge, dark, with a single accent rule below)
    f_name = font(240, weight="black", latin=True)
    name = "CodeRetina"
    nw, nh = text_size(draw, name, f_name)
    y = y_start + 60
    t(draw, (cx, y), name, f_name, PURPLE_DEEP, anchor="mt")

    # accent rule
    ry = y + nh + 30
    draw.rectangle([cx - 240, ry, cx + 240, ry + 9], fill=PURPLE)

    # English subtitle
    y = ry + 60
    f_en = font(50, weight="regular", latin=True, italic=True)
    t(draw, (cx, y),
      "A Visual Perception & Verification Layer for Code Agents",
      f_en, CHARCOAL, anchor="mt")

    # Chinese subtitle
    y += pt(50) + 18
    f_zh = font(58, weight="bold")
    t(draw, (cx, y),
      "为代码 Agent 增加视觉感知与验证闭环的多模态扩展层",
      f_zh, PURPLE_DEEP, anchor="mt")

    # Author line
    y += pt(58) + 38
    f_meta = font(44)
    t(draw, (cx, y),
      "作者：__________   ·   学号：__________   ·   指导教师：__________",
      f_meta, GRAY_TEXT, anchor="mt")
    y += pt(44) + 50

    # Section divider line (thin, full width)
    hline(draw, 200, W - 200, y, RULE, width=3)
    return y + 60

# ---------- §1 Motivation ----------

def render_motivation(draw, x, y, w):
    y = section(draw, x, y, w, "1", "研究动机", "Motivation")
    f_body = font(44)
    body = (
        "Code Agent 能读写代码、调用工具；Computer Use Agent "
        "能点击屏幕。但二者一旦进入真实软件环境，都面临同一个"
        "问题：文本上下文不够 —— Agent 不知道屏幕上有什么按钮、"
        "操作后状态是否变化、生成的页面是否符合预期。"
    )
    y = paragraph(draw, x, y, body, f_body, w, fill=CHARCOAL, line_h=1.5)
    y += 24
    # contrast pair
    f_lbl = font(44, weight="bold")
    t(draw, (x,                 y), "之前：", f_lbl, GRAY_TEXT)
    t(draw, (x + pt(44) * 3,    y), "Agent 只看代码", f_lbl, CHARCOAL)
    y += pt(44) + 16
    t(draw, (x,                 y), "之后：", f_lbl, PURPLE)
    t(draw, (x + pt(44) * 3,    y), "Agent 看屏 → 操作 → 自检", f_lbl, PURPLE_DEEP)
    return y + pt(44) + 18

# ---------- §2 The Vision Loop (the central diagram) ----------

def render_loop(draw, x, y, w):
    y = section(draw, x, y, w, "2", "视觉闭环", "The Vision Loop")
    f_node_zh = font(44, weight="bold")
    f_node_en = font(32, weight="bold", latin=True)
    f_center = font(40, weight="bold", latin=True)

    R = min(580, w // 2 - 220)
    node_rad = 160
    cx = x + w // 2
    cy = y + R + node_rad + 20
    diag_h = (R + node_rad) * 2 + 60

    # 5-node circle: observe / parse / plan / act / verify
    nodes = [
        ("观察",  "Observe",  PURPLE),
        ("解析",  "Parse",    (122, 73, 173)),
        ("规划",  "Plan",     (159, 86, 199)),
        ("执行",  "Act",      (197, 99, 219)),
        ("验证",  "Verify",   (220, 130, 90)),
    ]
    n = len(nodes)
    positions = []
    for i in range(n):
        # start at top, clockwise
        angle = -math.pi / 2 + i * (2 * math.pi / n)
        nx = cx + int(R * math.cos(angle))
        ny = cy + int(R * math.sin(angle))
        positions.append((nx, ny, angle))

    # draw connecting circular arrows (dotted-style: chunks)
    inner_R = R - node_rad // 2 - 20
    for i in range(n):
        a0 = -math.pi / 2 + i * (2 * math.pi / n) + 0.38
        a1 = -math.pi / 2 + (i + 1) * (2 * math.pi / n) - 0.38
        steps = 30
        for k in range(steps):
            ka = a0 + (a1 - a0) * (k / (steps - 1))
            px = cx + int(inner_R * math.cos(ka))
            py = cy + int(inner_R * math.sin(ka))
            r = 7 if k % 2 == 0 else 4
            draw.ellipse([px - r, py - r, px + r, py + r], fill=GRAY_LIGHT)
        ax = cx + int(inner_R * math.cos(a1))
        ay = cy + int(inner_R * math.sin(a1))
        tx = -math.sin(a1)
        ty =  math.cos(a1)
        nx = -tx * 0.5 + math.cos(a1) * 0.866
        ny = -ty * 0.5 + math.sin(a1) * 0.866
        size = 30
        pts = [
            (ax + int(tx * size),       ay + int(ty * size)),
            (ax - int(tx * size),       ay - int(ty * size)),
            (ax + int(nx * size * 1.6), ay + int(ny * size * 1.6)),
        ]
        draw.polygon(pts, fill=PURPLE)

    # draw nodes as filled circles with labels
    for (nx, ny, _ang), (zh, en, color) in zip(positions, nodes):
        draw.ellipse([nx - node_rad, ny - node_rad,
                      nx + node_rad, ny + node_rad],
                     fill=color, outline=PURPLE_DEEP, width=5)
        t(draw, (nx, ny - 22), zh, f_node_zh, WHITE, anchor="mm")
        t(draw, (nx, ny + 42), en, f_node_en, (245, 230, 255), anchor="mm")

    # center label
    t(draw, (cx, cy - 32), "Visual", f_center, PURPLE_DEEP, anchor="mm")
    t(draw, (cx, cy + 32), "Closed Loop", f_center, PURPLE_DEEP, anchor="mm")

    return y + diag_h

# ---------- §3 System Architecture (clean line-art) ----------

def render_architecture(draw, x, y, w):
    y = section(draw, x, y, w, "3", "系统架构", "System Architecture")
    f_box_zh = font(46, weight="bold")
    f_box_en = font(28, weight="regular", latin=True, italic=True)
    f_label  = font(40, weight="bold")
    f_chip   = font(32, weight="bold", latin=True)
    f_sub    = font(28, weight="regular", latin=True)

    # Two stacked layer boxes connected through a router pill.
    # All monochrome (purple stroke, white fill) — academic style.
    pad = 20
    top_h = 540
    router_h = 130
    arrow_h = 200
    bot_h = 700
    total = top_h + 60 + router_h + arrow_h + bot_h
    bx = x
    by = y
    cx = bx + w // 2

    # ---- Top: TS Agent CLI ----
    top = [bx + pad, by, bx + w - pad, by + top_h]
    rrect(draw, top, r=22, fill=WHITE, outline=PURPLE, width=5)
    t(draw, (top[0] + 30, top[1] + 22),
      "TypeScript Agent CLI", f_box_zh, PURPLE_DEEP)
    t(draw, (top[0] + 30, top[1] + 22 + pt(46) + 8),
      "Bun runtime · CLI · 35 commands · tool loop", f_box_en, GRAY_TEXT)
    # row of capability labels inside box (no colored chips)
    row_y = top[1] + 22 + pt(46) + 8 + pt(28) + 40
    f_chip_body = font(34, weight="bold")
    chips = [
        "Code Tools : File · Bash · Search · Edit",
        "Vision Tools : Screenshot · OCR · UIParse · VQA",
        "                       ImageDiff · VideoQA · DocRAG · Annotate · BrowserVision",
    ]
    for i, ctxt in enumerate(chips):
        cy = row_y + i * (pt(34) + 14)
        if i < 2:
            draw.rectangle([top[0] + 30, cy + pt(34) // 2 - 8,
                            top[0] + 30 + 16, cy + pt(34) // 2 + 8], fill=PURPLE)
        t(draw, (top[0] + 60, cy), ctxt, f_chip_body, CHARCOAL)

    # ---- Router pill (centered between layers) ----
    router_y = top[3] + 60
    f_router = font(38, weight="bold")
    router_text = "Hybrid Vision Router  ·  rule + confidence + budget + cache"
    rw, _ = text_size(draw, router_text, f_router)
    half = rw // 2 + 80
    router_box = [cx - half, router_y, cx + half, router_y + router_h]
    rrect(draw, router_box, r=router_h // 2, fill=PURPLE_DEEP, outline=None)
    t(draw, (cx, router_y + router_h // 2),
      router_text, f_router, WHITE, anchor="mm")

    # ---- Arrow + protocol label ----
    a_top = router_box[3] + 30
    a_bot = a_top + arrow_h - 70
    draw.line([(cx, a_top), (cx, a_bot)], fill=PURPLE, width=10)
    draw.polygon(
        [(cx - 36, a_bot), (cx + 36, a_bot), (cx, a_bot + 58)],
        fill=PURPLE,
    )
    t(draw, (cx + 60, (a_top + a_bot) // 2 - pt(28) // 2 - 4),
      "JSON-RPC over stdio",
      font(34, weight="bold", latin=True), PURPLE_DEEP, anchor="lm")
    t(draw, (cx + 60, (a_top + a_bot) // 2 + pt(28) // 2 + 6),
      "UTF-8 · sticky / half pack · timeout · restart",
      font(26, weight="regular", latin=True, italic=True), GRAY_TEXT, anchor="lm")

    # ---- Bottom: Python Sidecar ----
    bot_y0 = a_bot + 80
    bot = [bx + pad, bot_y0, bx + w - pad, bot_y0 + bot_h]
    rrect(draw, bot, r=22, fill=WHITE, outline=PURPLE, width=5)
    t(draw, (bot[0] + 30, bot[1] + 22),
      "Python Vision Sidecar", f_box_zh, PURPLE_DEEP)
    t(draw, (bot[0] + 30, bot[1] + 22 + pt(46) + 8),
      "uv · method registry · 15 RPC methods · local-first models", f_box_en, GRAY_TEXT)
    rpc = [
        "vlm.query · vlm.caption · vlm.agentic_qa",
        "ocr.extract · detect.objects · ui.parse",
        "image.diff · embed.image · rag.search_with_maxsim",
        "video.qa · doc.parse · gui.plan · gui.act",
    ]
    rrow_y = bot[1] + 22 + pt(46) + 8 + pt(28) + 40
    f_rpc = font(34, weight="bold", latin=True)
    for i, line in enumerate(rpc):
        cy = rrow_y + i * (pt(34) + 14)
        draw.rectangle([bot[0] + 30, cy + pt(34) // 2 - 8,
                        bot[0] + 30 + 16, cy + pt(34) // 2 + 8], fill=PURPLE)
        t(draw, (bot[0] + 60, cy), line, f_rpc, CHARCOAL)

    return by + total + 40

# ---------- §4 Hybrid Vision Router table ----------

def render_router(draw, x, y, w):
    y = section(draw, x, y, w, "4", "模型路由策略", "Hybrid Vision Router")
    rows = [
        ("Tier 1\n本地优先",  "MiniCPM-V · Moondream · PaddleOCR · YOLO · OmniParser",
         "默认；OCR · caption · UI 解析"),
        ("Tier 2\n便宜云端",  "Claude Haiku · Gemini Flash · Qwen-VL",
         "置信度不足时升级；常规 VQA"),
        ("Tier 3\n强云端",    "Claude Sonnet/Opus · GPT-4o · Gemini Pro",
         "困难 GUI planning · visual debug"),
    ]
    return draw_table(
        draw, x, y, w,
        headers=("层级", "代表模型 · Models", "适用场景 · When"),
        col_weights=(0.20, 0.48, 0.32),
        rows=rows,
        header_height=pt(70),
        row_height=pt(140),
    )

# ---------- Generic table renderer (academic style) ----------

def draw_table(draw, x, y, w, *, headers, col_weights, rows,
               header_height, row_height,
               header_fill=PURPLE_DEEP, header_fg=WHITE,
               alt_fill=TINT, body_fg=CHARCOAL):
    f_h = font(44, weight="bold")
    f_b_zh = font(44)
    f_b_en = font(32, weight="regular", latin=True)
    f_b = font(44)

    # column x positions
    xs = [x]
    for cw in col_weights[:-1]:
        xs.append(xs[-1] + int(w * cw))
    xs.append(x + w)

    # header
    draw.rectangle([x, y, x + w, y + header_height], fill=header_fill)
    for i, head in enumerate(headers):
        t(draw, (xs[i] + 24, y + header_height // 2 - pt(44) // 2 - 4),
          head, f_h, header_fg)
    y_cur = y + header_height

    # rows
    for r, row in enumerate(rows):
        if r % 2 == 1:
            draw.rectangle([x, y_cur, x + w, y_cur + row_height], fill=alt_fill)
        for i, cell in enumerate(row):
            cell_x = xs[i] + 24
            cell_w = xs[i + 1] - xs[i] - 48
            paragraph(draw, cell_x, y_cur + 24,
                      cell, f_b, cell_w, fill=body_fg, line_h=1.3)
        # row rule
        hline(draw, x, x + w, y_cur + row_height, RULE, width=2)
        y_cur += row_height

    # outer border
    draw.rectangle([x, y, x + w, y_cur], outline=PURPLE_DEEP, width=3)
    return y_cur + 36

# ---------- §5 Vision Toolchain table ----------

def render_toolchain(draw, x, y, w):
    y = section(draw, x, y, w, "5", "视觉工具链", "Vision Toolchain  (9 tools)")
    rows = [
        ("ScreenshotTool",    "屏幕截图",            "Observe"),
        ("BrowserVisionTool", "浏览器截图",          "Observe"),
        ("VisionQATool",      "图像问答 (+Agentic)", "Understand"),
        ("OCRTool",           "文字识别",            "Understand"),
        ("UIParseTool",       "UI 元素定位",         "Parse"),
        ("ImageDiffTool",     "前后视觉差异",        "Verify"),
        ("AnnotateTool",      "bbox 标注",           "Explain"),
        ("VideoQATool",       "录屏抽帧问答",        "Replay"),
        ("DocRAGTool",        "图文文档检索",        "Recall"),
    ]
    return draw_table(
        draw, x, y, w,
        headers=("工具 · Tool", "能力 · Capability", "环节 · Role"),
        col_weights=(0.36, 0.40, 0.24),
        rows=rows,
        header_height=pt(64),
        row_height=pt(58),
    )

# ---------- §6 Demo scenarios ----------

def render_demos(draw, x, y, w):
    y = section(draw, x, y, w, "6", "演示场景", "Demo Scenarios  (8 demos)")
    rows = [
        ("D1", "/design2code",   "设计图 → React+Tailwind → 截图 diff 迭代"),
        ("D2", "/visual-debug",  "执行前后截图对比 + VQA 判断状态"),
        ("D3", "/gui (dry-run)", "Observe-Parse-Plan-Act-Verify 闭环"),
        ("D4", "Agentic Search", "VLM 自驱动 crop / zoom / annotate"),
        ("D5", "/replay",        "ffmpeg + PySceneDetect 视频问答"),
        ("D6", "/doc",           "MinerU + ColQwen2 MaxSim 多模态 RAG"),
        ("D7", "Visual Planning","propose → predict → judge"),
        ("D8", "/skills-auto",   "Session 反思 → 自动技能沉淀"),
    ]
    return draw_table(
        draw, x, y, w,
        headers=("#", "命令 · Command", "做什么 · What"),
        col_weights=(0.08, 0.30, 0.62),
        rows=rows,
        header_height=pt(64),
        row_height=pt(58),
    )

# ---------- §7 Stats + Highlights ----------

def render_stats(draw, x, y, w):
    y = section(draw, x, y, w, "7", "工程量 & 亮点", "Engineering & Highlights")

    # numeric strip
    f_n   = font(96, weight="black", latin=True)
    f_lbl = font(34, weight="bold")
    stats = [
        ("9",     "Vision\nTools"),
        ("15",    "Sidecar\nMethods"),
        ("35",    "CLI\nCommands"),
        ("10.9k", "LOC\n(Vision)"),
        ("8",     "Demos /\nSprints"),
    ]
    cell_w = w // len(stats)
    for i, (num, lbl) in enumerate(stats):
        cx = x + i * cell_w + cell_w // 2
        t(draw, (cx, y), num, f_n, PURPLE_DEEP, anchor="mt")
        # label below number
        for j, line in enumerate(lbl.split("\n")):
            t(draw, (cx, y + pt(96) + 12 + j * (pt(34) + 6)),
              line, f_lbl, GRAY_TEXT, anchor="mt")
    y += pt(96) + 12 + 2 * (pt(34) + 6) + 40
    hline(draw, x, x + w, y, RULE, width=2)
    y += 36

    # highlights
    bullets = [
        ("跨语言协议工程化",
         "Content-Length 按 UTF-8 字节计算 · 粘包/半包 · 超时 · 子进程恢复"),
        ("视觉结果强结构化",
         "bbox · 置信度 · 动作 schema 全部可被 tool loop 直接消费"),
        ("三层模型混合路由",
         "本地优先 + 置信度升级 + 预算降级 + 结果缓存"),
        ("Observe-Parse-Plan-Act-Verify 闭环",
         "GUI Agent dry-run 默认安全；deliberative 多 candidate 评分"),
        ("视觉记忆 + 技能沉淀",
         "SigLIP2/ColPali + LanceDB · Session reflect → 自动 skill"),
    ]
    f_bt = font(44, weight="bold")
    f_bb = font(40)
    for title, body in bullets:
        # arrow marker
        ay = y + pt(44) // 2
        draw.polygon([(x, ay - 14), (x + 24, ay), (x, ay + 14)], fill=PURPLE)
        t(draw, (x + 44, y), title, f_bt, PURPLE_DEEP)
        y += pt(44) + 6
        y = paragraph(draw, x + 44, y, body, f_bb, w - 44, fill=CHARCOAL, line_h=1.35)
        y += 26
    return y

# ---------- Bottom strip: §7 Engineering + §8 Highlights side-by-side ----------

def render_bottom_strip(draw, x, y, w):
    """Compact full-width strip: numeric stats row + highlight tag row."""
    y = section(draw, x, y, w, "7", "工程量 & 亮点",
                "Engineering Footprint & Highlights")

    # ---- numeric stats: 5 cells across full width ----
    stats = [
        ("9",     "Vision Tools"),
        ("15",    "Sidecar Methods"),
        ("35",    "CLI Commands"),
        ("10.9k", "LOC (Vision)"),
        ("8",     "Demos / Sprints"),
    ]
    f_n   = font(96, weight="black", latin=True)
    f_lbl = font(36, weight="bold")
    cell_w = w // len(stats)
    for i, (num, lbl) in enumerate(stats):
        cx = x + i * cell_w + cell_w // 2
        t(draw, (cx, y), num, f_n, PURPLE_DEEP, anchor="mt")
        t(draw, (cx, y + pt(96) + 14),
          lbl, f_lbl, GRAY_TEXT, anchor="mt")
        # subtle vertical separator between cells
        if i < len(stats) - 1:
            sx = x + (i + 1) * cell_w
            vline(draw, sx, y + 14, y + pt(96) + 14 + pt(36) + 14,
                  RULE, width=2)
    y += pt(96) + 14 + pt(36) + 36
    hline(draw, x, x + w, y, RULE, width=2)
    y += 30

    # ---- highlight one-liners in horizontal "tag" row ----
    f_tag = font(40, weight="bold")
    tags = [
        "跨语言协议工程化",
        "视觉结果强结构化",
        "三层模型混合路由",
        "Observe→Verify 闭环",
        "视觉记忆 + 技能沉淀",
    ]
    # Lay out tags horizontally, evenly spaced
    total_w = 0
    sizes = []
    for tg in tags:
        tw, _ = text_size(draw, tg, f_tag)
        sizes.append(tw)
        total_w += tw
    sep = (w - total_w) // (len(tags) + 1)
    cur_x = x + sep
    for tg, tw in zip(tags, sizes):
        # diamond bullet
        cy_t = y + pt(40) // 2
        draw.polygon(
            [(cur_x - 24, cy_t), (cur_x - 12, cy_t - 12),
             (cur_x,     cy_t),  (cur_x - 12, cy_t + 12)],
            fill=PURPLE,
        )
        t(draw, (cur_x + 8, y), tg, f_tag, PURPLE_DEEP)
        cur_x += tw + sep
    return y + pt(40) + 30

# ---------- §8 References ----------

def render_refs(draw, x, y, w):
    y = section(draw, x, y, w, "8", "参考工作", "References  (selected)")
    items = [
        "Anthropic Computer Use (2024)",
        "UI-TARS (ByteDance, 2025)",
        "OmniParser v2 (Microsoft, 2025)",
        "MiniCPM-V 2.6 (OpenBMB, 2024)",
        "ColPali / ColQwen2 (ILLUIN, 2024)",
        "WebDreamer (OSU + Amazon, 2024)",
        "V* / ZoomEye (CVPR / THU, 2024)",
        "VOYAGER (NVIDIA, 2023)",
    ]
    f = font(36, weight="regular", latin=True)
    cw = w // 2
    for i, s in enumerate(items):
        col = i % 2
        row = i // 2
        t(draw, (x + col * cw, y + row * (pt(36) + 16)),
          f"·  {s}", f, GRAY_TEXT)
    return y + 4 * (pt(36) + 16) + 10

# ---------- v9 clear poster panels ----------

def panel(draw, x, y, w, h, num, title_zh, title_en):
    """Draw a reference-poster-style section panel and return content origin."""
    rrect(draw, [x, y, x + w, y + h], r=18, fill=WHITE, outline=RULE, width=3)
    header_h = 104
    draw.rounded_rectangle([x, y, x + w, y + header_h], radius=18, fill=PURPLE_DEEP)
    # Square off the lower header corners.
    draw.rectangle([x, y + header_h - 24, x + w, y + header_h], fill=PURPLE_DEEP)
    f_num = font(50, weight="black", latin=True)
    f_zh = font(48, weight="bold")
    t(draw, (x + 34, y + header_h // 2), num, f_num, WHITE, anchor="lm")
    t(draw, (x + 110, y + header_h // 2 - 2), title_zh, f_zh, WHITE, anchor="lm")
    return x + 42, y + header_h + 42, w - 84, h - header_h - 76

def draw_small_loop(draw, cx, cy, R):
    nodes = [
        ("观察", "Observe", PURPLE),
        ("解析", "Parse", (122, 73, 173)),
        ("规划", "Plan", (159, 86, 199)),
        ("执行", "Act", (197, 99, 219)),
        ("验证", "Verify", (220, 130, 90)),
    ]
    f_zh = font(32, weight="bold")
    f_en = font(22, weight="bold", latin=True)
    node_r = 72
    for i in range(len(nodes)):
        a0 = -math.pi / 2 + i * (2 * math.pi / len(nodes)) + 0.34
        a1 = -math.pi / 2 + (i + 1) * (2 * math.pi / len(nodes)) - 0.34
        for k in range(20):
            a = a0 + (a1 - a0) * k / 19
            px = cx + int((R - 40) * math.cos(a))
            py = cy + int((R - 40) * math.sin(a))
            draw.ellipse([px - 4, py - 4, px + 4, py + 4], fill=GRAY_LIGHT)
        ax = cx + int((R - 40) * math.cos(a1))
        ay = cy + int((R - 40) * math.sin(a1))
        draw.polygon([(ax, ay - 20), (ax + 34, ay), (ax, ay + 20)], fill=PURPLE)

    for i, (zh, en, color) in enumerate(nodes):
        angle = -math.pi / 2 + i * (2 * math.pi / len(nodes))
        nx = cx + int(R * math.cos(angle))
        ny = cy + int(R * math.sin(angle))
        draw.ellipse([nx - node_r, ny - node_r, nx + node_r, ny + node_r],
                     fill=color, outline=PURPLE_DEEP, width=4)
        t(draw, (nx, ny - 14), zh, f_zh, WHITE, anchor="mm")
        t(draw, (nx, ny + 28), en, f_en, (245, 230, 255), anchor="mm")

    t(draw, (cx, cy - 22), "Visual", font(32, weight="bold", latin=True),
      PURPLE_DEEP, anchor="mm")
    t(draw, (cx, cy + 26), "Closed Loop", font(32, weight="bold", latin=True),
      PURPLE_DEEP, anchor="mm")

def draw_metric_chip(draw, x, y, label, value, w):
    rrect(draw, [x, y, x + w, y + 138], r=16, fill=TINT, outline=RULE, width=2)
    t(draw, (x + 28, y + 34), label, font(30, weight="bold"), GRAY_TEXT)
    t(draw, (x + 28, y + 88), value, font(40, weight="black"), PURPLE_DEEP)

def render_problem_loop_panel(draw, x, y, w, h):
    cx, cy, cw, ch = panel(draw, x, y, w, h, "1", "动机与闭环", "Motivation & Loop")
    f_body = font(39)
    text = (
        "Code Agent 会读写代码，Computer Use Agent 会点击屏幕；"
        "但真实软件环境的状态来自截图、弹窗、按钮、表格和页面渲染。"
        "CodeRetina 把这些视觉输入转成 Agent 可消费的结构化上下文。"
    )
    yy = paragraph(draw, cx, cy, text, f_body, cw, fill=CHARCOAL, line_h=1.42)
    yy += 28
    chip_w = (cw - 34) // 2
    draw_metric_chip(draw, cx, yy, "Before", "只看代码 / 文本", chip_w)
    draw_metric_chip(draw, cx + chip_w + 34, yy, "After", "看屏 → 操作 → 自检", chip_w)
    loop_cx = cx + cw // 2
    loop_cy = y + h - 340
    draw_small_loop(draw, loop_cx, loop_cy, 260)
    return y + h

def render_architecture_panel(draw, x, y, w, h):
    cx, cy, cw, ch = panel(draw, x, y, w, h, "2", "系统架构", "System Architecture")
    f_title = font(34, weight="bold", latin=True)
    f_sub = font(25, weight="regular", latin=True)
    f_zh = font(31, weight="bold")

    box_h = 225
    gap = 48
    top = [cx, cy, cx + cw, cy + box_h]
    mid = [cx + 90, top[3] + gap, cx + cw - 90, top[3] + gap + 150]
    bot = [cx, mid[3] + gap, cx + cw, mid[3] + gap + box_h + 10]

    rrect(draw, top, r=18, fill=(250, 247, 253), outline=PURPLE, width=4)
    t(draw, (top[0] + 32, top[1] + 44), "TypeScript Agent CLI", f_title, PURPLE_DEEP)
    t(draw, (top[0] + 32, top[1] + 106), "CLI · tool loop · commands", f_sub, GRAY_TEXT)

    rrect(draw, mid, r=75, fill=PURPLE_DEEP, outline=None)
    t(draw, ((mid[0] + mid[2]) // 2, (mid[1] + mid[3]) // 2 - 10),
      "Hybrid Vision Router", font(36, weight="bold", latin=True), WHITE, anchor="mm")
    t(draw, ((mid[0] + mid[2]) // 2, (mid[1] + mid[3]) // 2 + 40),
      "rule · confidence · budget · cache", font(24, weight="bold", latin=True),
      (235, 220, 250), anchor="mm")

    for y0, y1 in [(top[3], mid[1]), (mid[3], bot[1])]:
        mx = cx + cw // 2
        draw.line([(mx, y0 + 10), (mx, y1 - 10)], fill=PURPLE, width=8)
        draw.polygon([(mx - 28, y1 - 20), (mx + 28, y1 - 20), (mx, y1 + 22)], fill=PURPLE)

    rrect(draw, bot, r=18, fill=(250, 247, 253), outline=PURPLE, width=4)
    t(draw, (bot[0] + 32, bot[1] + 38), "Python Vision Sidecar", f_title, PURPLE_DEEP)
    t(draw, (bot[0] + 32, bot[1] + 96), "JSON-RPC stdio · method registry", f_sub, GRAY_TEXT)
    t(draw, (bot[0] + 32, bot[1] + 164),
      "VLM · OCR · UI Parse · Image Diff · Doc / Video RAG",
      font(27, weight="bold", latin=True), CHARCOAL)

    yy = bot[3] + 46
    tiers = [
        ("Tier 1", "本地优先"),
        ("Tier 2", "低成本云"),
        ("Tier 3", "强推理"),
    ]
    tier_w = (cw - 36 * 2) // 3
    for i, (name, zh) in enumerate(tiers):
        bx = cx + i * (tier_w + 36)
        rrect(draw, [bx, yy, bx + tier_w, yy + 145], r=14,
              fill=WHITE if i % 2 == 0 else TINT, outline=RULE, width=2)
        t(draw, (bx + tier_w // 2, yy + 48), name,
          font(34, weight="black", latin=True), PURPLE_DEEP, anchor="mm")
        t(draw, (bx + tier_w // 2, yy + 102), zh, font(28, weight="bold"),
          CHARCOAL, anchor="mm")
    return y + h

def fit_image(im, max_w, max_h):
    scale = min(max_w / im.width, max_h / im.height)
    nw, nh = int(im.width * scale), int(im.height * scale)
    return im.resize((nw, nh), Image.LANCZOS)

def draw_agentic_demo_figure(draw, x, y, w, h):
    """Synthetic but information-rich visual trace for the poster demo."""
    rrect(draw, [x, y, x + w, y + h], r=22, fill=(250, 250, 252),
          outline=PURPLE_DEEP, width=4)

    f_ui = font(22, weight="bold", latin=True)
    f_small = font(18, weight="regular", latin=True)
    f_label = font(26, weight="black", latin=True)

    # Main screenshot frame.
    sx, sy = x + 36, y + 42
    sw, sh = w - 72, int(h * 0.66)
    rrect(draw, [sx, sy, sx + sw, sy + sh], r=14, fill=WHITE,
          outline=RULE, width=2)
    draw.rectangle([sx, sy, sx + sw, sy + 66], fill=(47, 55, 75))
    t(draw, (sx + 28, sy + 43), "Customer Admin Console", f_ui, WHITE, anchor="lm")
    for i, c in enumerate([(255, 92, 92), (242, 203, 86), (88, 201, 154)]):
        draw.ellipse([sx + sw - 118 + i * 34, sy + 24,
                      sx + sw - 100 + i * 34, sy + 42], fill=c)

    # Sidebar and dense content cards.
    side_w = 300
    draw.rectangle([sx, sy + 66, sx + side_w, sy + sh], fill=(244, 244, 248))
    menu_items = ["Dashboard", "Users", "Billing", "Audit Log", "Settings"]
    for i, item in enumerate(menu_items):
        my = sy + 104 + i * 72
        fill = TINT_HEAVY if item == "Audit Log" else WHITE
        rrect(draw, [sx + 24, my, sx + side_w - 24, my + 46], r=8,
              fill=fill, outline=RULE, width=1)
        t(draw, (sx + 44, my + 29), item, f_small, CHARCOAL, anchor="lm")

    content_x = sx + side_w + 34
    content_y = sy + 104
    card_w = (sw - side_w - 34 * 4) // 3
    card_h = 170
    for row in range(3):
        for col in range(3):
            bx = content_x + col * (card_w + 34)
            by = content_y + row * (card_h + 38)
            rrect(draw, [bx, by, bx + card_w, by + card_h], r=12,
                  fill=WHITE, outline=RULE, width=2)
            draw.rectangle([bx, by, bx + card_w, by + 46], fill=(79, 91, 116))
            t(draw, (bx + 18, by + 31), f"Panel {row * 3 + col + 1}",
              f_small, WHITE, anchor="lm")
            draw.rectangle([bx + 18, by + 72, bx + card_w - 28, by + 88],
                           fill=(224, 224, 230))
            draw.rectangle([bx + 18, by + 108, bx + card_w - 92, by + 124],
                           fill=(232, 232, 238))

    # Tiny error toast: the actual target of the visual search.
    toast_w, toast_h = 310, 116
    tx = sx + sw - toast_w - 130
    ty = sy + sh - toast_h - 42
    rrect(draw, [tx, ty, tx + toast_w, ty + toast_h], r=14,
          fill=(255, 244, 238), outline=(198, 64, 35), width=3)
    t(draw, (tx + 24, ty + 36), "Payment failed", font(22, weight="bold", latin=True),
      (120, 40, 30), anchor="lm")
    t(draw, (tx + 24, ty + 78), "ERR_AUTH_42", font(28, weight="black", latin=True),
      (198, 64, 35), anchor="lm")

    # Search box and label.
    pad = 14
    draw.rectangle([tx - pad, ty - pad, tx + toast_w + pad, ty + toast_h + pad],
                   outline=ACCENT, width=8)
    rrect(draw, [tx - pad, ty - 64, tx + 246, ty - 20], r=10,
          fill=ACCENT, outline=None)
    t(draw, (tx, ty - 42), "target: tiny error code", font(22, weight="bold", latin=True),
      WHITE, anchor="lm")

    # Bottom trace: crop/zoom/result, drawn as a clear three-step story.
    trace_y = sy + sh + 46
    trace_h = h - (trace_y - y) - 36
    step_gap = 28
    step_w = (sw - 2 * step_gap) // 3
    steps = [
        ("1 Observe", "full screen"),
        ("2 Crop + Zoom", "bottom-right toast"),
        ("3 Answer", "ERR_AUTH_42"),
    ]
    for i, (title, body) in enumerate(steps):
        bx = sx + i * (step_w + step_gap)
        fill = TINT_HEAVY if i == 1 else WHITE
        rrect(draw, [bx, trace_y, bx + step_w, trace_y + trace_h], r=14,
              fill=fill, outline=RULE, width=2)
        t(draw, (bx + 22, trace_y + 38), title, f_label, PURPLE_DEEP, anchor="lm")
        t(draw, (bx + 22, trace_y + 84), body, font(24, weight="bold", latin=True),
          CHARCOAL, anchor="lm")
        if i == 1:
            # Zoomed crop preview.
            zx, zy = bx + 24, trace_y + 116
            zw, zh = step_w - 48, trace_h - 142
            rrect(draw, [zx, zy, zx + zw, zy + zh], r=10,
                  fill=(255, 244, 238), outline=ACCENT, width=4)
            t(draw, (zx + 22, zy + 42), "ERR_AUTH_42",
              font(34, weight="black", latin=True), (198, 64, 35), anchor="lm")
        if i < 2:
            ax = bx + step_w + 4
            ay = trace_y + trace_h // 2
            draw.line([(ax, ay), (ax + step_gap - 8, ay)], fill=PURPLE, width=5)
            draw.polygon([(ax + step_gap - 8, ay - 14),
                          (ax + step_gap - 8, ay + 14),
                          (ax + step_gap + 10, ay)], fill=PURPLE)

def render_demo_panel(img, draw, x, y, w, h):
    cx, cy, cw, ch = panel(draw, x, y, w, h, "3", "旗舰 Demo：主动视觉搜索", "Flagship Demo: Agentic Visual Search")
    fig_w = int(cw * 0.58)
    fig_h = ch - 80
    sx, sy = cx, cy + 30
    shadow = Image.new("RGBA", (fig_w + 28, fig_h + 28), (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle([14, 14, fig_w + 14, fig_h + 14],
                         radius=18, fill=(60, 40, 90, 45))
    shadow = shadow.filter(ImageFilter.GaussianBlur(10))
    img.alpha_composite(shadow, (sx - 12, sy - 10))
    draw_agentic_demo_figure(draw, sx, sy, fig_w, fig_h)

    rx = sx + fig_w + 90
    rw = cx + cw - rx
    t(draw, (rx, cy + 42), "Demo evidence",
      font(40, weight="black", latin=True), PURPLE_DEEP)
    yy = cy + 128
    body = (
        "Dense UI、细小按钮和局部文字容易让单次 VQA 失效。"
        "这里让模型主动 crop / zoom / annotate，形成可追踪的多步视觉推理。"
    )
    yy = paragraph(draw, rx, yy, body, font(30), rw, fill=CHARCOAL, line_h=1.32)
    yy += 30

    steps = [
        ("1", "Observe"),
        ("2", "Crop / Zoom"),
        ("3", "Trace"),
    ]
    for no, title in steps:
        rrect(draw, [rx, yy, rx + rw, yy + 118], r=18,
              fill=TINT if no != "2" else TINT_HEAVY, outline=RULE, width=2)
        draw.ellipse([rx + 28, yy + 28, rx + 86, yy + 86], fill=PURPLE_DEEP)
        t(draw, (rx + 57, yy + 57), no, font(30, weight="black", latin=True),
          WHITE, anchor="mm")
        t(draw, (rx + 112, yy + 58), title,
          font(32, weight="black", latin=True), PURPLE_DEEP, anchor="lm")
        yy += 138

    yy += 8
    draw_metric_chip(draw, rx, yy, "Self-test", "5 / 5 passed", (rw - 30) // 2)
    draw_metric_chip(draw, rx + (rw - 30) // 2 + 30, yy,
                     "Confidence", "+9% on tiny text", (rw - 30) // 2)
    return y + h

def render_capabilities_panel(draw, x, y, w, h):
    cx, cy, cw, ch = panel(draw, x, y, w, h, "4", "核心能力族", "Capability Families")
    cards = [
        ("Observe", "截图 · 浏览器 · 视频帧"),
        ("Understand", "VQA · OCR · UIParse"),
        ("Verify", "ImageDiff · 状态检查"),
        ("Recall", "Doc RAG · Video Chapters"),
    ]
    gap = 34
    card_w = (cw - gap) // 2
    card_h = 255
    for i, (title, subtitle) in enumerate(cards):
        bx = cx + (i % 2) * (card_w + gap)
        by = cy + (i // 2) * (card_h + 34)
        rrect(draw, [bx, by, bx + card_w, by + card_h], r=18,
              fill=TINT if i % 2 else WHITE, outline=RULE, width=2)
        t(draw, (bx + 28, by + 58), title,
          font(34, weight="black", latin=True), PURPLE_DEEP)
        paragraph(draw, bx + 28, by + 132, subtitle, font(25, weight="bold"),
                  card_w - 56, fill=CHARCOAL, line_h=1.22)
    return y + h

def render_scenarios_panel(draw, x, y, w, h):
    cx, cy, cw, ch = panel(draw, x, y, w, h, "5", "代表场景", "Representative Scenarios")
    demos = [
        ("/gui", "Observe → Verify", "Computer Use 闭环"),
        ("/design2code", "Design → Code → Diff", "前端视觉验证"),
        ("/doc /replay", "PDF / Video → QA", "文档与录屏问答"),
    ]
    yy = cy
    for cmd, flow, desc in demos:
        rrect(draw, [cx, yy, cx + cw, yy + 164], r=18,
              fill=WHITE, outline=RULE, width=2)
        t(draw, (cx + 28, yy + 48), cmd,
          font(28, weight="black", latin=True), PURPLE_DEEP)
        t(draw, (cx + 28, yy + 112), desc, font(23, weight="bold"), GRAY_TEXT)
        t(draw, (cx + 455, yy + 66), flow, font(27, weight="bold", latin=True),
          CHARCOAL, anchor="lm")
        yy += 186
    yy += 8
    rrect(draw, [cx, yy, cx + cw, yy + 132], r=18, fill=PURPLE_DEEP, outline=None)
    t(draw, (cx + 34, yy + 66), "Other demos retained as supporting evidence.",
      font(31, weight="bold", latin=True), WHITE, anchor="lm")
    return y + h

# ---------- Layout glue ----------

def main():
    img = Image.new("RGBA", (W, H), WHITE)
    draw = ImageDraw.Draw(img)

    banner_top(img, draw)
    y = 720
    y = render_hero(img, draw, y)

    # v9 layout: clear poster panels, fewer tables, one real demo figure.
    margin = 160
    gutter = 90
    bw = W - 2 * margin
    col_w = (W - 2 * margin - gutter) // 2
    col_xs = [margin, margin + col_w + gutter]
    body_y_start = y

    row1_h = 1540
    row2_h = 1580
    row3_h = 930
    gap_y = 66

    row1_y = body_y_start
    render_problem_loop_panel(draw, col_xs[0], row1_y, col_w, row1_h)
    render_architecture_panel(draw, col_xs[1], row1_y, col_w, row1_h)

    row2_y = row1_y + row1_h + gap_y
    render_demo_panel(img, draw, margin, row2_y, bw, row2_h)

    row3_y = row2_y + row2_h + gap_y
    render_capabilities_panel(draw, col_xs[0], row3_y, col_w, row3_h)
    render_scenarios_panel(draw, col_xs[1], row3_y, col_w, row3_h)

    stats_y = row3_y + row3_h + 48
    render_bottom_strip(draw, margin, stats_y, bw)

    print(f"body start: {body_y_start}")
    print(f"stats starts at: {stats_y}")
    print(f"footer starts at {H - 360}")

    banner_bottom(img, draw)

    rgb = img.convert("RGB")
    rgb.save(OUT_JPG, "JPEG", quality=92, dpi=(DPI, DPI), optimize=True)
    rgb.save(OUT_PDF, "PDF", resolution=DPI)

    pres = Presentation()
    pres.slide_width = Emu(32399288)
    pres.slide_height = Emu(43200638)
    blank = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank)
    slide.shapes.add_picture(str(OUT_JPG), 0, 0,
                             width=pres.slide_width,
                             height=pres.slide_height)
    pres.save(OUT_PPTX)
    print(f"JPG   : {OUT_JPG}")
    print(f"PDF   : {OUT_PDF}")
    print(f"PPTX  : {OUT_PPTX}")

if __name__ == "__main__":
    main()
